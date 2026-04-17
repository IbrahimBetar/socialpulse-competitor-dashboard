import io
import json
import time
import traceback
from datetime import datetime, timedelta, date

import pandas as pd
import plotly.express as px
import streamlit as st
import yt_dlp
from yt_dlp import YoutubeDL
from pptx import Presentation
from pptx.util import Inches, Pt

from datetime import datetime, date

def format_platform_name(raw_input, platform):
    """Cleans URLs into readable names with platform logos."""
    if platform == "YouTube":
        if "@" in raw_input:
            clean_name = raw_input.split("@")[-1].split("/")[0]
        else:
            clean_name = raw_input.split("/")[-1] if "/" in raw_input else raw_input
        return f"🔴 YouTube: @{clean_name}"
    elif platform == "TikTok":
        clean_name = raw_input.replace("@", "").strip()
        return f"🎵 TikTok: @{clean_name}"
    return raw_input


def is_within_range(post_date_str, start_date, end_date):
    """Helper function to check if a string date falls within the selected date range."""
    try:
        # Handle mixed formats like '2026-04-13T12:00:00Z' or '2026-04-13'
        if "T" in post_date_str or len(post_date_str) > 10:
            post_date = datetime.strptime(post_date_str[:10], "%Y-%m-%d").date()
        else:
            post_date = datetime.strptime(post_date_str, "%Y-%m-%d").date()
            
        return start_date <= post_date <= end_date
    except (ValueError, TypeError):
        # If the date cannot be parsed, default to True so we don't accidentally drop valid data
        return True

START_DATE = datetime(2025, 1, 1)
END_DATE = datetime(2026, 4, 11)


# Create a silent logger to prevent yt-dlp from crashing when it hits an error
class SilentLogger(object):
    def debug(self, msg): pass
    def warning(self, msg): pass
    def error(self, msg): pass


def scrape_youtube(channel_url, start_date, end_date):
    total_likes = 0
    total_comments = 0
    total_views = 0
    timeline_data = []
    
    if not channel_url.endswith('/videos'):
        channel_url = f"{channel_url.rstrip('/')}/videos"
        
    yt_opts = {
        'quiet': True,
        'extract_flat': False,
        'skip_download': True,
        'js_runtimes': {'node': {}},
        'playlistend': 200,
        'ignoreerrors': True
    }
    
    try:
        with YoutubeDL(yt_opts) as ydl:
            info = ydl.extract_info(channel_url, download=False)
            
            if not info or 'entries' not in info:
                return 0, 0, 0, pd.DataFrame()
                
            for video in info['entries']:
                if not video:
                    continue
                
                upload_date_str = video.get('upload_date')
                if not upload_date_str:
                    continue
                    
                video_date = datetime.strptime(upload_date_str, '%Y%m%d')
                
                if video_date > end_date:
                    continue
                if video_date < start_date:
                    break
                    
                views = video.get('view_count', 0) or 0
                likes = video.get('like_count', 0) or 0
                comments = video.get('comment_count', 0) or 0
                
                total_views += views
                total_likes += likes
                total_comments += comments
                
                timeline_data.append({
                    'date': video_date.strftime('%Y-%m-%d'),
                    'platform': 'YouTube',
                    'engagement': likes + comments,  # Views are impressions, not engagement
                    'views': views,
                    'likes': likes,
                    'comments': comments,
                    'shares': 0, # Default to 0 for YT compatibility
                    'saves': 0,  # Default to 0 for YT compatibility
                    'title': video.get('title', 'Unknown')
                })
                
        df = pd.DataFrame(timeline_data)
        return total_likes, total_comments, total_views, df
        
    except Exception as e:
        st.error(f"YouTube Error: {str(e)}")
        print(traceback.format_exc())
        return 0, 0, 0, pd.DataFrame()


def scrape_tiktok(handle):
    if not handle:
        return None

    profile_url = f"https://www.tiktok.com/@{handle}"
    ydl_opts = {
        "quiet": True,
        "skip_download": True,
        "extract_flat": False,
        "js_runtimes": {"node": {}},
        "playlistend": 200,
        "ignoreerrors": True,
        "extractor_args": {"tiktok": ["api_hostname=api22-normal-c-useast2a.tiktokv.com"]}
    }
    result = {
        "platform": "TikTok",
        "total_likes": 0,
        "total_comments": 0,
        "total_views": 0,
        "total_shares": 0,
        "total_saves": 0,
        "posts": [],
    }

    with YoutubeDL(ydl_opts) as ydl:
        info = ydl.extract_info(profile_url, download=False)
    time.sleep(3)

    entries = info.get("entries") or [info]
    for entry in entries:
        if entry is None:
            continue

        upload_date = entry.get("upload_date")
        if not upload_date:
            continue

        if not is_within_range(upload_date, START_DATE.date(), END_DATE.date()):
            continue

        published = datetime.strptime(upload_date, "%Y%m%d")

        views = entry.get("play_count", 0) or entry.get("view_count", 0) or 0
        likes = entry.get("digg_count", 0) or entry.get("like_count", 0) or 0
        comments = entry.get("comment_count", 0) or 0
        
        # Safely try to get shares/saves (different extractors use different keys)
        shares = entry.get('repost_count') or entry.get('share_count', 0)
        saves = entry.get('save_count') or entry.get('bookmark_count', 0)

        result["total_likes"] += likes
        result["total_comments"] += comments
        result["total_views"] += views
        result["total_shares"] += shares
        result["total_saves"] += saves
        result["posts"].append(
            {
                "date": published.date(),
                "engagement": likes + comments + shares + saves,  # True engagement excludes views
                "views": views,
                "likes": likes,
                "comments": comments,
                "shares": shares,
                "saves": saves,
            }
        )

    return result


def build_metrics(platform_result):
    if not platform_result:
        return [0, 0, 0, 0, 0]
    return [
        platform_result.get("total_likes", 0),
        platform_result.get("total_comments", 0),
        platform_result.get("total_views", 0),
        platform_result.get("total_shares", 0),
        platform_result.get("total_saves", 0),
    ]


def generate_ppt_report(all_results, start_date, end_date):
    """Generate a PowerPoint presentation report from brand data."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Master Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Competitor Video Performance Report"
    subtitle.text = f"Analyzed between {start_date} and {end_date}\nGenerated by SocialPulse"

    for brand_name, platforms_data in all_results.items():
        if not platforms_data:
            continue

        comp_timeline_df = build_timeline(platforms_data)
        if comp_timeline_df.empty:
            continue

        for p_data in platforms_data:
            comp_timeline_df.loc[comp_timeline_df["platform"] == p_data["platform"], "platform"] = p_data.get(
                "platform_display", p_data["platform"]
            )

        mask = (comp_timeline_df["date"] >= start_date) & (comp_timeline_df["date"] <= end_date)
        filtered_df = comp_timeline_df.loc[mask]

        if filtered_df.empty:
            continue

        # Brand Summary Slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = f"Performance Summary: {brand_name}"

        total_views = filtered_df["views"].sum() if "views" in filtered_df.columns else 0
        total_likes = filtered_df["likes"].sum() if "likes" in filtered_df.columns else 0
        total_comments = filtered_df["comments"].sum() if "comments" in filtered_df.columns else 0
        total_engagement = filtered_df["engagement"].sum() if "engagement" in filtered_df.columns else 0

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        tf.text = f"Total Posts Analyzed: {len(filtered_df)}"
        tf.add_paragraph().text = f"Total Combined Views: {total_views:,.0f}"
        tf.add_paragraph().text = f"Total Combined Likes: {total_likes:,.0f}"
        tf.add_paragraph().text = f"Total Combined Comments: {total_comments:,.0f}"
        tf.add_paragraph().text = f"Total Combined Engagement: {total_engagement:,.0f}"

        # Top Posts Table Slide for this Brand
        table_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(table_slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = f"Top 5 Performing Posts: {brand_name}"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        top_posts = filtered_df.sort_values(by="views", ascending=False).head(5)

        rows = len(top_posts) + 1
        cols = 4
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(4.5)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(4.0)

        headers = ["Date", "Platform", "Views", "URL"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True

        for row_idx, (_, row) in enumerate(top_posts.iterrows(), start=1):
            table.cell(row_idx, 0).text = str(row["date"])
            table.cell(row_idx, 1).text = str(row["platform"])
            table.cell(row_idx, 2).text = f"{row.get('views', 0):,.0f}"
            table.cell(row_idx, 3).text = str(row.get("url", "N/A"))

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io


def display_summary_metrics(platform_results):
    if not platform_results:
        return

    # Dynamically calculate totals if metadata is missing
    for platform_result in platform_results:
        if not platform_result:
            continue
        if platform_result.get("total_views", 0) == 0 and "posts" in platform_result:
            platform_result["total_views"] = sum(
                post.get("views", 0)
                for post in platform_result["posts"]
                if isinstance(post.get("views", 0), (int, float))
            )
            platform_result["total_likes"] = sum(
                post.get("likes", 0)
                for post in platform_result["posts"]
                if isinstance(post.get("likes", 0), (int, float))
            )
            platform_result["total_comments"] = sum(
                post.get("comments", 0)
                for post in platform_result["posts"]
                if isinstance(post.get("comments", 0), (int, float))
            )
            platform_result["total_shares"] = sum(
                post.get("shares", 0)
                for post in platform_result["posts"]
                if isinstance(post.get("shares", 0), (int, float))
            )
            platform_result["total_saves"] = sum(
                post.get("saves", 0)
                for post in platform_result["posts"]
                if isinstance(post.get("saves", 0), (int, float))
            )

    for platform_result in platform_results:
        if not platform_result:
            continue
        platform = platform_result["platform"]
        likes, comments, views, shares, saves = build_metrics(platform_result)
        cols = st.columns(5)
        cols[0].metric(f"{platform} Likes", f"{likes:,}")
        cols[1].metric(f"{platform} Comments", f"{comments:,}")
        cols[2].metric(f"{platform} Shares", f"{shares:,}")
        cols[3].metric(f"{platform} Saves", f"{saves:,}")
        cols[4].metric(f"{platform} Views", f"{views:,}")


def build_timeline(results):
    rows = []
    for platform_result in results:
        if not platform_result:
            continue
        platform = platform_result.get("platform")
        for post in platform_result.get("posts", []):
            rows.append(
                {
                    "date": post["date"],
                    "engagement": post.get("engagement", 0),
                    "views": post.get("views", 0),
                    "likes": post.get("likes", 0),
                    "comments": post.get("comments", 0),
                    "shares": post.get("shares", 0),
                    "saves": post.get("saves", 0),
                    "platform": platform,
                }
            )
    if not rows:
        return pd.DataFrame(columns=["date", "engagement", "views", "likes", "comments", "shares", "saves", "platform"])

    df = pd.DataFrame(rows)
    df = df.groupby(["date", "platform"], as_index=False).sum()
    df["date"] = pd.to_datetime(df["date"]).dt.date
    return df.sort_values("date", ascending=False)


def main():
    st.set_page_config(page_title="Competitor Analysis Dashboard", layout="wide")
    st.title("Competitor Analysis Dashboard")
    st.write("Track engagement across YouTube and TikTok by Brand.")

    if "all_results" not in st.session_state:
        st.session_state.all_results = {}
    if "run_analysis" not in st.session_state:
        st.session_state.run_analysis = False

    st.sidebar.subheader("🏢 Competitor Brands")
    brands_input = []

    for i in range(1, 4):
        with st.sidebar.expander(f"Brand {i}", expanded=(i == 1)):
            b_name = st.text_input(f"Brand Name", key=f"name_{i}", placeholder="e.g., Nike")
            b_yt = st.text_input(f"YouTube URL", key=f"yt_{i}", placeholder="https://youtube.com/@nike")
            b_tt = st.text_input(f"TikTok Handle", key=f"tt_{i}", placeholder="nike")

            if b_name and (b_yt or b_tt):
                brands_input.append({"name": b_name, "youtube": b_yt, "tiktok": b_tt})

    col1, col2 = st.sidebar.columns(2)

    if col1.button("▶ Run Analysis"):
        st.session_state.run_analysis = True

    if col2.button("🗑 Clear Data"):
        st.session_state.all_results = {}
        st.session_state.run_analysis = False
        st.rerun()

    st.sidebar.subheader("📅 Date Range Filter")
    today = date.today()
    default_start = today - timedelta(days=30)

    date_selection = st.sidebar.date_input(
        "Analyze posts between:",
        value=(default_start, today),
        max_value=today
    )

    if getattr(st.session_state, "run_analysis", False):
        if not brands_input:
            st.warning("Please configure at least one brand.")
            st.stop()

        for brand in brands_input:
            brand_name = brand["name"]

            if brand_name not in st.session_state.all_results:
                st.session_state.all_results[brand_name] = []

                if brand.get("youtube"):
                    try:
                        clean_yt = format_platform_name(brand["youtube"], "YouTube")
                        with st.spinner(f"Scraping {clean_yt}..."):
                            total_likes, total_comments, total_views, yt_df = scrape_youtube(
                                brand["youtube"], START_DATE, END_DATE
                            )
                            yt_posts = []
                            if not yt_df.empty:
                                yt_posts = [
                                    {
                                        "date": row["date"],
                                        "engagement": row["engagement"],
                                        "platform": "YouTube",
                                        "views": row.get("views", 0),
                                        "likes": row.get("likes", 0),
                                        "comments": row.get("comments", 0),
                                        "shares": row.get("shares", 0),
                                        "saves": row.get("saves", 0),
                                        "title": row.get("title", "Unknown"),
                                    }
                                    for _, row in yt_df.iterrows()
                                ]

                            yt_data = {
                                "platform": "YouTube",
                                "platform_display": clean_yt,
                                "total_likes": total_likes,
                                "total_comments": total_comments,
                                "total_views": total_views,
                                "total_shares": 0,  # YouTube doesn't expose shares/saves publicly
                                "total_saves": 0,   # YouTube doesn't expose shares/saves publicly
                                "posts": yt_posts,
                            }
                            st.session_state.all_results[brand_name].append(yt_data)
                            st.success(f"✅ {clean_yt} loaded")
                    except Exception as exc:
                        st.error(f"YouTube error: {exc}")

                if brand.get("tiktok"):
                    try:
                        clean_tt = format_platform_name(brand["tiktok"], "TikTok")
                        with st.spinner(f"Scraping {clean_tt}..."):
                            tt_data = scrape_tiktok(brand["tiktok"])
                            if tt_data and tt_data.get("posts"):
                                tt_data["platform_display"] = clean_tt
                                st.session_state.all_results[brand_name].append(tt_data)
                                st.success(f"✅ {clean_tt} loaded")
                            else:
                                st.warning(f"No TikTok data found for {brand['tiktok']}.")
                    except Exception as exc:
                        st.error(f"TikTok error: {exc}")

        st.session_state.run_analysis = False

    if len(date_selection) == 2:
        start_date, end_date = date_selection
    else:
        start_date = end_date = date_selection

    if not st.session_state.all_results:
        st.info("Configure a brand and click '▶ Run Analysis' to start.")
    else:
        comparison_rows = []
        for comp_name, platforms_data in st.session_state.all_results.items():
            total_likes = 0
            total_comments = 0
            total_views = 0
            total_engagement = 0
            total_posts = 0
            platforms = sorted({result.get("platform") for result in platforms_data if result.get("platform")})

            for result in platforms_data:
                for post in result.get("posts", []):
                    try:
                        post_date = (
                            datetime.strptime(post["date"][:10], "%Y-%m-%d").date()
                            if isinstance(post["date"], str)
                            else post["date"]
                        )
                    except (TypeError, ValueError):
                        continue

                    if not (start_date <= post_date <= end_date):
                        continue

                    likes = post.get("likes", 0)
                    comments = post.get("comments", 0)
                    views = post.get("views", 0)
                    engagement = post.get("engagement", 0)

                    total_likes += likes
                    total_comments += comments
                    total_views += views
                    total_engagement += engagement
                    total_posts += 1

            comparison_rows.append(
                {
                    "competitor": comp_name,
                    "platforms": ", ".join(platforms),
                    "posts": total_posts,
                    "likes": total_likes,
                    "comments": total_comments,
                    "views": total_views,
                    "engagement": total_engagement,
                }
            )

        if comparison_rows:
            comparison_df = pd.DataFrame(comparison_rows)
            comparison_df = comparison_df.sort_values(["engagement", "views"], ascending=False)
            st.subheader("Competitor Comparison Summary")
            total_likes = 0
            total_comments = 0
            total_views = 0
            total_engagement = 0
            total_posts = 0
            platforms = sorted({result.get("platform") for result in platforms_data if result.get("platform")})

            for result in platforms_data:
                for post in result.get("posts", []):
                    try:
                        post_date = (
                            datetime.strptime(post["date"][:10], "%Y-%m-%d").date()
                            if isinstance(post["date"], str)
                            else post["date"]
                        )
                    except (TypeError, ValueError):
                        continue

                    if not (start_date <= post_date <= end_date):
                        continue

                    likes = post.get("likes", 0)
                    comments = post.get("comments", 0)
                    views = post.get("views", 0)
                    engagement = post.get("engagement", 0)

                    total_likes += likes
                    total_comments += comments
                    total_views += views
                    total_engagement += engagement
                    total_posts += 1

            comparison_rows.append({
                "competitor": comp_name,
                "platforms": ", ".join(platforms),
                "posts": total_posts,
                "likes": total_likes,
                "comments": total_comments,
                "views": total_views,
                "engagement": total_engagement,
            })

        if comparison_rows:
            comparison_df = pd.DataFrame(comparison_rows)
            comparison_df = comparison_df.sort_values(["engagement", "views"], ascending=False)
            st.subheader("Competitor Comparison Summary")
            st.dataframe(comparison_df, use_container_width=True)
            try:
                fig = px.bar(
                    comparison_df,
                    x="competitor",
                    y=["engagement", "views"],
                    title="Competitor Comparison: Engagement vs Views",
                    barmode="group",
                    text_auto=True,
                )
                fig.update_layout(xaxis_title="Competitor", yaxis_title="Total Count")
                st.plotly_chart(fig, use_container_width=True)
            except Exception:
                pass

        for brand_name, platforms_data in st.session_state.all_results.items():
            if not platforms_data:
                continue

            st.header(f"🏢 Brand Analysis: {brand_name}")

            for result in platforms_data:
                if result and result.get("posts"):
                    filtered_posts = []
                    for post in result.get("posts", []):
                        try:
                            if isinstance(post["date"], str):
                                if "T" in post["date"] or len(post["date"]) > 10:
                                    post_date = datetime.strptime(post["date"][:10], "%Y-%m-%d").date()
                                else:
                                    post_date = datetime.strptime(post["date"], "%Y-%m-%d").date()
                            else:
                                post_date = post["date"]

                            if start_date <= post_date <= end_date:
                                filtered_posts.append(post)
                        except (ValueError, TypeError):
                            filtered_posts.append(post)

                    result["posts"] = filtered_posts
                    result["total_likes"] = sum(post.get("likes", 0) for post in filtered_posts)
                    result["total_comments"] = sum(post.get("comments", 0) for post in filtered_posts)
                    result["total_views"] = sum(post.get("views", 0) for post in filtered_posts)

            comp_timeline_df = build_timeline(platforms_data)
            if not comp_timeline_df.empty:
                mask = (comp_timeline_df["date"] >= start_date) & (comp_timeline_df["date"] <= end_date)
                filtered_df = comp_timeline_df.loc[mask]
                for p_data in platforms_data:
                    filtered_df.loc[filtered_df["platform"] == p_data["platform"], "platform"] = p_data.get(
                        "platform_display", p_data["platform"]
                    )
            else:
                filtered_df = pd.DataFrame()

            st.subheader("Summary Metrics")
            display_summary_metrics(platforms_data)

            comp_raw_rows = []
            for result in platforms_data:
                platform_display = result.get("platform_display", result.get("platform", "Unknown"))
                for post in result.get("posts", []):
                    comp_raw_rows.append(
                        {
                            "date": post["date"].strftime("%Y-%m-%d")
                            if isinstance(post["date"], date)
                            else post["date"],
                            "platform": platform_display,
                            "engagement": post.get("engagement", 0),
                            "views": post.get("views", 0),
                            "likes": post.get("likes", 0),
                            "comments": post.get("comments", 0),
                        }
                    )

            comp_raw_df = pd.DataFrame(comp_raw_rows)
            if not comp_raw_df.empty:
                comp_platform_metrics = comp_raw_df.groupby("platform", as_index=False).agg(
                    {
                        "likes": "sum",
                        "comments": "sum",
                        "views": "sum",
                    }
                )
                st.subheader("Summary Table")
                st.dataframe(comp_platform_metrics, use_container_width=True)
            else:
                st.info(f"No summary table data available for {brand_name}.")

            st.subheader("Engagement Over Time")
            if filtered_df.empty:
                st.info(f"No timeline data available for {brand_name}.")
            else:
                chart_df = filtered_df.melt(
                    id_vars=["date", "platform"],
                    value_vars=["engagement", "views"],
                    var_name="metric",
                    value_name="value",
                )
                fig = px.line(
                    chart_df,
                    x="date",
                    y="value",
                    color="metric",
                    facet_col="platform",
                    title=f"Engagement and Views Over Time for {brand_name}",
                    markers=True,
                    color_discrete_map={"engagement": "#1f77b4", "views": "#ff7f0e"},
                )
                fig.update_layout(xaxis_title="Date", yaxis_title="Count", legend_title_text="Metric")
                st.plotly_chart(fig, use_container_width=True)

            st.subheader("Raw Post Timeline")
            if filtered_df.empty:
                st.info(f"No timeline data available for {brand_name}.")
            else:
                st.dataframe(filtered_df, use_container_width=True)

            # --- Monthly Performance Comparison for THIS Brand ---
            st.subheader(f"📅 Monthly Performance: {brand_name}")
            
            if not filtered_df.empty:
                monthly_source_df = filtered_df.copy()
            
                # Ensure required columns exist
                if 'likes' not in monthly_source_df.columns:
                    monthly_source_df['likes'] = 0
                if 'comments' not in monthly_source_df.columns:
                    monthly_source_df['comments'] = 0
                if 'views' not in monthly_source_df.columns:
                    monthly_source_df['views'] = 0
                if 'shares' not in monthly_source_df.columns:
                    monthly_source_df['shares'] = 0
                if 'saves' not in monthly_source_df.columns:
                    monthly_source_df['saves'] = 0
                if 'platform' not in monthly_source_df.columns:
                    monthly_source_df['platform'] = 'Unknown'
            
                # Ensure date is datetime
                monthly_source_df['date'] = pd.to_datetime(monthly_source_df['date'], errors='coerce')
                monthly_source_df = monthly_source_df.dropna(subset=['date'])
            
                if not monthly_source_df.empty:
                    monthly_source_df['month'] = monthly_source_df['date'].dt.to_period('M').astype(str)
                    # Engagement = Likes + Comments + Shares + Saves ONLY (No Views!)
                    monthly_source_df['engagement'] = (
                        monthly_source_df['likes'].fillna(0) + 
                        monthly_source_df['comments'].fillna(0) +
                        monthly_source_df.get('shares', pd.Series(0, index=monthly_source_df.index)).fillna(0) +
                        monthly_source_df.get('saves', pd.Series(0, index=monthly_source_df.index)).fillna(0)
                    )
                    monthly_source_df['posts'] = 1
            
                    # Group by platform and month
                    monthly_df = monthly_source_df.groupby(['platform', 'month'], as_index=False).agg({
                        'posts': 'sum',
                        'likes': 'sum',
                        'comments': 'sum',
                        'shares': 'sum',
                        'saves': 'sum',
                        'views': 'sum',
                        'engagement': 'sum'
                    })
            
                    # Sort and calculate MoM deltas per platform
                    monthly_df = monthly_df.sort_values(['platform', 'month'])
                    monthly_df['engagement_mom_delta'] = monthly_df.groupby('platform')['engagement'].diff().fillna(0)
                    monthly_df['views_mom_delta'] = monthly_df.groupby('platform')['views'].diff().fillna(0)
            
                    # Optional formatting
                    display_monthly_df = monthly_df.copy()
                    format_cols = ['posts', 'likes', 'comments', 'shares', 'saves', 'views', 'engagement', 'engagement_mom_delta', 'views_mom_delta']
                    for col in format_cols:
                        if col in display_monthly_df.columns:
                            display_monthly_df[col] = display_monthly_df[col].apply(lambda x: f"{int(x):,}")
            
                    st.dataframe(display_monthly_df, use_container_width=True, hide_index=True)
                else:
                    st.info(f'Not enough valid dated data to calculate monthly performance for {brand_name}.')
            else:
                st.info(f'Not enough data to calculate monthly performance for {brand_name}.')

            st.divider()

        st.sidebar.markdown("---")
        st.sidebar.subheader("📥 Export Data")
        ppt_file = generate_ppt_report(st.session_state.all_results, start_date, end_date)

        st.sidebar.download_button(
            label="📊 Download Full PowerPoint Report",
            data=ppt_file,
            file_name=f"SocialPulse_Analysis_{start_date}_to_{end_date}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
st.markdown("---")
st.caption("Contact: [elbetar2001@gmail.com](mailto:elbetar2001@gmail.com)")
if __name__ == "__main__":
    main()
